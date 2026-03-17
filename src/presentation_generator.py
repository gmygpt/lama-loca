"""
Генератор презентаций PPTX
"""
import os
import re
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import config


class PresentationGenerator:

    def __init__(self):
        os.makedirs(config.OUTPUT_DIR, exist_ok=True)

    def _sanitize(self, name: str) -> str:
        name = re.sub(r'[^\w\s\-]', '', name)
        name = re.sub(r'\s+', '_', name.strip())
        return name[:80]

    def _parse_slides(self, text: str) -> list:
        slides = []
        current = None

        for line in text.split('\n'):
            line = line.strip()
            if not line:
                continue

            m = re.match(
                r'^(?:СЛАЙД|Слайд|слайд|SLIDE|Slide)\s*(\d+)\s*[:\.\-]\s*(.*)', line
            )
            if m:
                if current:
                    slides.append(current)
                current = {"title": m.group(2).strip(), "points": []}
                continue

            if current is not None:
                cleaned = re.sub(r'^[\-\*•]\s*', '', line)
                cleaned = re.sub(r'^\d+[\.\)]\s*', '', cleaned)
                if cleaned:
                    current["points"].append(cleaned)

        if current:
            slides.append(current)
        return slides

    def generate(self, text: str, topic: str) -> str:
        slides_data = self._parse_slides(text)

        if not slides_data:
            slides_data = self._fallback(text, topic)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        first_title = slides_data[0]["title"] if slides_data else topic
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = first_title
        for p in sl.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        if sl.placeholders[1]:
            sl.placeholders[1].text = f"Подготовлено: {datetime.now().strftime('%d.%m.%Y')}"

        # Content slides
        start = 1 if slides_data and not slides_data[0]["points"] else 0
        for sd in slides_data[start:]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sd["title"]
            for p in slide.shapes.title.text_frame.paragraphs:
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, point in enumerate(sd["points"][:6]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = point
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x44)
                    p.space_after = Pt(8)

        filename = f"презентация_{self._sanitize(topic)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(config.OUTPUT_DIR, filename)
        prs.save(filepath)
        return filepath

    def _fallback(self, text: str, topic: str) -> list:
        slides = [{"title": topic, "points": []}]
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        for para in paragraphs[:10]:
            lines = [l.strip() for l in para.split('\n') if l.strip()]
            if lines:
                slides.append({"title": lines[0][:80], "points": lines[1:] or [para[:200]]})
        slides.append({"title": "Заключение", "points": ["Спасибо за внимание!", f"Тема: {topic}"]})
        return slides
